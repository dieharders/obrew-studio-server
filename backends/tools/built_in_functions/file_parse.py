"""File parse tool for extracting text from documents using existing loaders."""
from pathlib import Path
from typing import Dict, Any, Literal
from pydantic import BaseModel, Field


class Params(BaseModel):
    """Extract full text content from documents (PDF, DOCX, PPTX, XLSX, RTF, CSV, XML, HTML, MD)."""

    file_path: str = Field(
        ...,
        description="The path to the document file to parse.",
    )
    output_format: Literal["markdown", "text"] = Field(
        default="text",
        description="Output format: 'markdown' adds basic formatting, 'text' is plain text.",
    )

    model_config = {
        "json_schema_extra": {
            "examples": [
                {
                    "file_path": "/documents/report.pdf",
                    "output_format": "text",
                },
                {
                    "file_path": "/documents/presentation.pptx",
                    "output_format": "text",
                },
            ]
        }
    }


# Supported document extensions mapped to their parser
SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".rtf": "rtf",
    ".csv": "csv",
    ".xml": "xml",
    ".html": "html",
    ".htm": "html",
    ".md": "text",
    ".markdown": "text",
    ".txt": "text",
    ".json": "text",
}


def _parse_pdf(file_path: Path) -> tuple[str, int]:
    """Parse PDF using PyMuPDF (fitz)."""
    import fitz
    doc = fitz.open(str(file_path))
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    page_count = len(doc)
    doc.close()
    return "\n".join(text_parts), page_count


def _parse_docx(file_path: Path) -> str:
    """Parse DOCX using python-docx."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise ValueError(
            "python-docx is required for .docx files. Install with: pip install python-docx"
        )
    doc = DocxDocument(str(file_path))
    text_parts = [para.text for para in doc.paragraphs]
    return "\n".join(text_parts)


def _parse_pptx(file_path: Path) -> str:
    """Parse PPTX using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError(
            "python-pptx is required for .pptx files. Install with: pip install python-pptx"
        )
    prs = Presentation(str(file_path))
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
    return "\n".join(text_parts)


def _parse_xlsx(file_path: Path) -> str:
    """Parse XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ValueError(
            "openpyxl is required for .xlsx files. Install with: pip install openpyxl"
        )
    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    text_parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text_parts.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            row_text = ", ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                text_parts.append(row_text)
    wb.close()
    return "\n".join(text_parts)


def _parse_rtf(file_path: Path) -> str:
    """Parse RTF using striprtf."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ValueError(
            "striprtf is required for .rtf files. Install with: pip install striprtf"
        )
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        rtf_content = f.read()
    return rtf_to_text(rtf_content)


def _parse_csv(file_path: Path) -> str:
    """Parse CSV file."""
    import csv
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = [", ".join(row) for row in reader]
    return "\n".join(rows)


def _parse_xml(file_path: Path) -> str:
    """Parse XML file."""
    import xml.etree.ElementTree as ET
    tree = ET.parse(str(file_path))
    root = tree.getroot()

    def get_text(element):
        text_parts = []
        if element.text:
            text_parts.append(element.text.strip())
        for child in element:
            text_parts.extend(get_text(child))
            if child.tail:
                text_parts.append(child.tail.strip())
        return text_parts

    return " ".join(filter(None, get_text(root)))


def _parse_html(file_path: Path) -> str:
    """Parse HTML file - extract text content."""
    import re
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    # Remove script and style tags
    content = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL | re.IGNORECASE)
    content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL | re.IGNORECASE)
    # Remove HTML tags
    content = re.sub(r'<[^>]+>', ' ', content)
    # Clean up whitespace
    content = re.sub(r'\s+', ' ', content)
    return content.strip()


def _parse_text(file_path: Path) -> str:
    """Parse plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


async def main(**kwargs) -> Dict[str, Any]:
    """
    Parse a document and extract its text content.
    Returns dict with: content, format, file_path, page_count (if available)
    """
    file_path_str = kwargs.get("file_path")
    output_format = kwargs.get("output_format", "text")

    if not file_path_str:
        raise ValueError("file_path is required")

    file_path = Path(file_path_str)

    if not file_path.exists():
        raise ValueError(f"File does not exist: {file_path_str}")
    if not file_path.is_file():
        raise ValueError(f"Path is not a file: {file_path_str}")

    extension = file_path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file format: {extension}. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    parser_type = SUPPORTED_EXTENSIONS[extension]
    page_count = None
    method = parser_type

    try:
        if parser_type == "pdf":
            content, page_count = _parse_pdf(file_path)
            method = "fitz"
        elif parser_type == "docx":
            content = _parse_docx(file_path)
            method = "python-docx"
        elif parser_type == "pptx":
            content = _parse_pptx(file_path)
            method = "python-pptx"
        elif parser_type == "xlsx":
            content = _parse_xlsx(file_path)
            method = "openpyxl"
        elif parser_type == "rtf":
            content = _parse_rtf(file_path)
            method = "striprtf"
        elif parser_type == "csv":
            content = _parse_csv(file_path)
            method = "csv"
        elif parser_type == "xml":
            content = _parse_xml(file_path)
            method = "xml"
        elif parser_type == "html":
            content = _parse_html(file_path)
            method = "html"
        else:  # text
            content = _parse_text(file_path)
            method = "direct_read"

        # For markdown output format, add basic structure
        if output_format == "markdown" and parser_type not in ["text"]:
            content = f"# {file_path.name}\n\n{content}"

        return {
            "content": content,
            "format": output_format,
            "file_path": str(file_path.resolve()),
            "page_count": page_count,
            "method": method,
        }

    except Exception as e:
        raise ValueError(f"Failed to parse document: {e}")
