import os
import csv
import json
import xml.etree.ElementTree as ET
from typing import List, Optional
from dotenv import load_dotenv
import fitz  # PyMuPDF
from core import classes, common
from core.document import Document

# Optional llama-parse for cloud PDF parsing (requires API key)
# Currently dep is removed.
# from llama_parse import LlamaParse

###########
# METHODS #
###########


def create_source_document(text: str, source_id: str, metadata: dict):
    """Create a Document with metadata."""
    document_metadata = {
        **metadata,
        "sourceId": source_id,
    }
    source_doc = Document(
        id_=source_id,
        text=text,
        metadata=document_metadata,
    )
    return source_doc


def set_ignored_metadata(source_document: Document, ignore_metadata: dict):
    """Set metadata keys to ignore during embedding/LLM processing."""
    source_document.excluded_llm_metadata_keys.extend(["sourceId", "embedder"])
    source_document.excluded_embed_metadata_keys.extend(["sourceId", "embedder"])
    source_document.metadata.update(ignore_metadata)
    source_document.excluded_llm_metadata_keys.extend(list(ignore_metadata.keys()))
    source_document.excluded_embed_metadata_keys.extend(list(ignore_metadata.keys()))
    return source_document


###########
# LOADERS #
###########


def simple_file_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load simple text files (txt, md, json, etc.)."""
    documents = []
    for path in sources:
        with open(path, "r", encoding="utf-8") as f:
            text = f.read()

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc, ignore_metadata=source_metadata
        )
        documents.append(source_doc)
    return documents


def simple_pdf_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
    make_metadata: bool = True,
) -> List[Document]:
    """Load PDF files using PyMuPDF (fitz)."""
    document_results: List[Document] = []
    for path in sources:
        doc = fitz.open(path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())

        metadata = (
            {
                "total_pages": len(doc),
                **source_metadata,
            }
            if make_metadata
            else source_metadata
        )

        doc.close()

        source_doc = create_source_document(
            text="".join(text_parts),
            source_id=source_id,
            metadata=metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


def ms_doc_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load Microsoft Word documents (.docx)."""
    document_results: List[Document] = []

    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise ImportError(
            "python-docx is required for .docx files. Install with: pip install python-docx"
        )

    for path in sources:
        doc = DocxDocument(path)
        text_parts = [para.text for para in doc.paragraphs]

        source_doc = create_source_document(
            text="\n".join(text_parts),
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


def rtf_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load RTF files."""
    document_results: List[Document] = []

    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ImportError(
            "striprtf is required for .rtf files. Install with: pip install striprtf"
        )

    for path in sources:
        with open(path, "r", encoding="utf-8") as f:
            rtf_content = f.read()
        text = rtf_to_text(rtf_content)

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


def csv_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load CSV files."""
    document_results: List[Document] = []
    for path in sources:
        with open(path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = [", ".join(row) for row in reader]
            text = "\n".join(rows)

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


def xml_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load XML files."""
    document_results: List[Document] = []
    for path in sources:
        tree = ET.parse(path)
        root = tree.getroot()

        # Extract all text content from XML
        def get_text(element):
            text_parts = []
            if element.text:
                text_parts.append(element.text.strip())
            for child in element:
                text_parts.extend(get_text(child))
                if child.tail:
                    text_parts.append(child.tail.strip())
            return text_parts

        text = " ".join(filter(None, get_text(root)))

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


def pptx_slides_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load PowerPoint files (.pptx)."""
    document_results: List[Document] = []

    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for .pptx files. Install with: pip install python-pptx"
        )

    for path in sources:
        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)

        source_doc = create_source_document(
            text="\n".join(text_parts),
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


def simple_image_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load image files. Returns empty text (no OCR)."""
    document_results: List[Document] = []
    for path in sources:
        # Without OCR, we just note that an image exists
        filename = os.path.basename(path)
        text = f"[Image: {filename}]"

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


async def vision_image_loader(
    app: classes.FastAPIApp,
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load images using vision model for transcription.

    Uses the loaded vision model to generate text descriptions of images.
    Falls back to simple_image_loader if vision model is not available or fails.
    """
    document_results: List[Document] = []

    # Check if vision model is loaded
    if not hasattr(app.state, "vision_llm") or not app.state.vision_llm:
        print(
            f"{common.PRNT_EMBED} No vision model loaded. "
            "Falling back to placeholder text.",
            flush=True,
        )
        return simple_image_loader(sources, source_id, source_metadata)

    vision_llm = app.state.vision_llm

    for path in sources:
        filename = os.path.basename(path)
        try:
            print(
                f"{common.PRNT_EMBED} Transcribing image with vision model: {filename}",
                flush=True,
            )

            # Run vision inference (non-streaming for embedding)
            prompt = (
                "Describe this image in detail. Include any visible text, "
                "objects, people, scenes, colors, and relevant context. "
                "If there is text in the image, transcribe it exactly."
            )

            response_gen = await vision_llm.vision_completion(
                prompt=prompt,
                image_paths=[path],
                request=None,  # No request context needed for embedding
                stream=False,
            )

            # Collect response from generator
            text = ""
            async for chunk in response_gen:
                if isinstance(chunk, dict) and "content" in chunk:
                    text = chunk["content"]
                elif isinstance(chunk, str):
                    # Parse JSON string response
                    try:
                        parsed = json.loads(chunk)
                        if "content" in parsed:
                            text = parsed["content"]
                    except json.JSONDecodeError:
                        pass

            if not text:
                raise Exception("Empty response from vision model")

            print(
                f"{common.PRNT_EMBED} Vision transcription complete for: {filename}",
                flush=True,
            )

        except Exception as e:
            print(
                f"{common.PRNT_EMBED} Vision transcription failed for {filename}: {e}. "
                "Falling back to placeholder.",
                flush=True,
            )
            text = f"[Image: {filename}]"

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)

    return document_results


def simple_audio_video_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Load audio/video files. Returns empty text (no transcription)."""
    document_results: List[Document] = []
    for path in sources:
        # Without transcription, we just note that a media file exists
        filename = os.path.basename(path)
        text = f"[Media: {filename}]"

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)
    return document_results


# Stub functions for advanced loaders
def smart_pdf_loader():
    return []


def scientific_pdf_loader():
    return []


def image_vision_loader():
    return []


def unstructured_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Fallback loader - reads file as text."""
    return simple_file_loader(sources, source_id, source_metadata)


async def llama_parse_loader(
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """
    LlamaParse cloud service loader for advanced PDF parsing.

    Requires:
    - llama-parse package: pip install llama-parse
    - LLAMA_CLOUD_API_KEY environment variable

    Falls back to simple_pdf_loader as this is not currently being used.
    """

    print(
        f"{common.PRNT_EMBED} llama-parse not installed. "
        "Falling back to simple PDF loader. "
        "Install with: pip install llama-parse",
        flush=True,
    )
    return simple_pdf_loader(sources, source_id, source_metadata)

    load_dotenv()
    llama_parse_api_key = os.getenv("LLAMA_CLOUD_API_KEY")
    document_results: List[Document] = []

    if not llama_parse_api_key:
        print(f"{common.PRNT_EMBED} LLAMA_CLOUD_API_KEY not set.", flush=True)
        return document_results

    parser = LlamaParse(
        api_key=llama_parse_api_key,
        result_type="markdown",
        num_workers=8,
        verbose=True,
        language="en",
    )

    for path in sources:
        try:
            results = await parser.aload_data(path)
            for result in results:
                source_doc = create_source_document(
                    text=result.text if hasattr(result, "text") else str(result),
                    source_id=source_id,
                    metadata=source_metadata,
                )
                source_doc = set_ignored_metadata(
                    source_document=source_doc,
                    ignore_metadata=source_metadata,
                )
                document_results.append(source_doc)
        except Exception as e:
            print(
                f"{common.PRNT_EMBED} LlamaParse failed for {path}: {e}. "
                "Falling back to simple PDF loader.",
                flush=True,
            )
            fallback_docs = simple_pdf_loader([path], source_id, source_metadata)
            document_results.extend(fallback_docs)

    return document_results


def jina_reader_loader(
    app: classes.FastAPIApp,
    sources: List[str],
    source_id: str,
    source_metadata: dict,
) -> List[Document]:
    """Free web reader using Jina AI."""
    document_results: List[Document] = []
    for path in sources:
        req_url = f"https://r.jina.ai/{path}"
        headers = {
            "Accept": "text/event-stream",
            "Content-Type": "application/octet-stream",
        }
        client = app.state.requests_client
        text = ""
        with client.stream(method="GET", url=req_url, headers=headers) as res:
            res.raise_for_status()
            if res.status_code == 200:
                res.read()
                text = res.text
            else:
                raise Exception("Something went wrong reading data.")

        source_doc = create_source_document(
            text=text,
            source_id=source_id,
            metadata=source_metadata,
        )
        source_doc = set_ignored_metadata(
            source_document=source_doc,
            ignore_metadata=source_metadata,
        )
        document_results.append(source_doc)

    return document_results


async def documents_from_sources(
    app: classes.FastAPIApp,
    sources: List[str],
    source_id: str,
    source_metadata: dict,
    parsing_method: Optional[classes.FILE_LOADER_SOLUTIONS] = None,
) -> List[Document]:
    """Read source files and build document nodes with metadata."""
    print(f"{common.PRNT_EMBED} Reading files...", flush=True)
    documents = []
    for source in sources:
        filename = os.path.basename(source)
        file_extension = common.get_file_extension_from_path(filename).lower()
        payload = dict(
            sources=[source],
            source_id=source_id,
            source_metadata=source_metadata,
        )
        # Use loader solution based on file type
        match file_extension:
            case "mdx" | "md" | "json" | "txt":
                documents = simple_file_loader(**payload)
            case "doc" | "docx":
                documents = ms_doc_loader(**payload)
            case "rtf":
                documents = rtf_loader(**payload)
            case "csv":
                documents = csv_loader(**payload)
            case "xml":
                documents = xml_loader(**payload)
            case "pptx":
                documents = pptx_slides_loader(**payload)
            case "png" | "jpg" | "jpeg" | "gif":
                # Use vision model for transcription if available
                documents = await vision_image_loader(app=app, **payload)
            case "mp4" | "mp3":
                documents = simple_audio_video_loader(**payload)
            case "pdf":
                match (parsing_method):
                    case classes.FILE_LOADER_SOLUTIONS.LLAMA_PARSE.value:
                        documents = await llama_parse_loader(**payload)
                    case _:
                        documents = simple_pdf_loader(**payload)
            case _:
                is_url = source[:4] == "http"
                if (
                    is_url
                    and parsing_method == classes.FILE_LOADER_SOLUTIONS.READER.value
                ):
                    documents = jina_reader_loader(
                        app=app,
                        **payload,
                    )
                else:
                    raise Exception(
                        f"The supplied file/url is not currently supported: {source}"
                    )
    return documents
